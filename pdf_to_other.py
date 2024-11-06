# -- coding: utf-8 --
from pdf2docx import Converter
import os
from pptx import Presentation
from openpyxl import Workbook

def pdf_to_word(pdf_file, word_file):
    # 使用 pdf2docx 将 PDF 转换为 Word 文件
    cv = Converter(pdf_file)
    cv.convert(word_file, start=0, end=None)
    cv.close()

def pdf_to_ppt(pdf_file, ppt_file):
    # 简单示例，将每个 PDF 页面内容转换为 PPT 的一页
    prs = Presentation()
    # 假设 PDF 已经有某种文本数据
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "PDF 转 PPT 示例"
    prs.save(ppt_file)

def pdf_to_excel(pdf_file, excel_file):
    # 简单示例，将 PDF 的内容转换为 Excel
    wb = Workbook()
    ws = wb.active
    ws.append(["PDF 转 Excel 示例"])
    wb.save(excel_file)
